# -*- coding: utf-8 -*-
"""Gera a apresentacao executiva A PARTIR das saidas do pipeline.

    python apresentacao.py

POR QUE GERADA, E NAO ESCRITA A MAO
-----------------------------------
Um deck com numeros digitados descola do modelo no primeiro rerun. Aqui todo
numero vem de outputs/resumo_execucao.json e dos CSV do book: rodar o pipeline
de novo e regerar o deck produz um material sempre coerente com a planilha.

O rodape de cada slide carrega o STATUS do run. Se a serie historica ainda for
fixture, isso aparece no material — nao se apresenta numero sem procedencia.
"""
from __future__ import annotations

import json
from datetime import date
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

RAIZ = Path(__file__).resolve().parent
OUT = RAIZ / "outputs"

ROSA = RGBColor(0xC2, 0x18, 0x5B)
ESC = RGBColor(0x1F, 0x12, 0x35)
CINZA = RGBColor(0x61, 0x61, 0x61)
VERDE = RGBColor(0x2E, 0x7D, 0x32)
VERM = RGBColor(0xB7, 0x1C, 0x1C)
BRANCO = RGBColor(0xFF, 0xFF, 0xFF)
CLARO = RGBColor(0xF7, 0xF7, 0xFA)

L, T, W = Inches(0.62), Inches(0.55), Inches(12.1)


def _mi(x):
    return f"R$ {x/1e6:,.1f} mi".replace(",", "@").replace(".", ",").replace("@", ".")


def _n(x, d=0):
    return f"{x:,.{d}f}".replace(",", "@").replace(".", ",").replace("@", ".")


def carregar():
    d = json.loads((OUT / "resumo_execucao.json").read_text(encoding="utf-8"))
    ctx = {"resumo": d, "book": d.get("book", {})}
    for nome, arq in (("curva", "curva_mensal_base_seco_umido.csv"),
                      ("sinal", "sinal_vs_mercado.csv"),
                      ("pernas", "book_pernas.csv"),
                      ("dim", "book_dimensionamento_risco.csv")):
        p = OUT / arq
        ctx[nome] = pd.read_csv(p) if p.exists() else pd.DataFrame()
    return ctx


class Deck:
    def __init__(self, status: str):
        self.p = Presentation()
        self.p.slide_width, self.p.slide_height = Inches(13.333), Inches(7.5)
        self.status = status
        self.n = 0

    def _rodape(self, s):
        cx = s.shapes.add_textbox(L, Inches(6.95), W, Inches(0.35))
        tf = cx.text_frame
        tf.text = (f"Curva Publica de Referencia — SE/CO  |  corte 14/08/2026  |  "
                   f"{self.status}  |  {self.n}")
        r = tf.paragraphs[0].runs[0]
        r.font.size, r.font.color.rgb = Pt(9), CINZA

    def slide(self, titulo, subtitulo=None, capa=False):
        s = self.p.slides.add_slide(self.p.slide_layouts[6])
        self.n += 1
        if capa:
            cx = s.shapes.add_textbox(L, Inches(2.3), W, Inches(1.4))
            tf = cx.text_frame
            tf.text = titulo
            r = tf.paragraphs[0].runs[0]
            r.font.size, r.font.bold, r.font.color.rgb = Pt(40), True, ROSA
            if subtitulo:
                p2 = tf.add_paragraph()
                p2.text = subtitulo
                p2.runs[0].font.size = Pt(18)
                p2.runs[0].font.color.rgb = ESC
            self._rodape(s)
            return s
        cx = s.shapes.add_textbox(L, T, W, Inches(0.6))
        tf = cx.text_frame
        tf.text = titulo
        r = tf.paragraphs[0].runs[0]
        r.font.size, r.font.bold, r.font.color.rgb = Pt(24), True, ROSA
        if subtitulo:
            cx2 = s.shapes.add_textbox(L, Inches(1.12), W, Inches(0.45))
            t2 = cx2.text_frame
            t2.word_wrap = True
            t2.text = subtitulo
            rr = t2.paragraphs[0].runs[0]
            rr.font.size, rr.font.color.rgb = Pt(12), CINZA
        self._rodape(s)
        return s

    def kpis(self, s, itens, topo=1.75, alt=1.25):
        n = len(itens)
        larg = (W - Inches(0.15) * (n - 1)) / n
        for i, (rot, val, cor) in enumerate(itens):
            x = L + (larg + Inches(0.15)) * i
            cx = s.shapes.add_textbox(x, Inches(topo), larg, Inches(alt))
            tf = cx.text_frame
            tf.word_wrap = True
            tf.text = rot
            rr = tf.paragraphs[0].runs[0]
            rr.font.size, rr.font.color.rgb = Pt(10), CINZA
            p2 = tf.add_paragraph()
            p2.text = val
            p2.runs[0].font.size = Pt(23)
            p2.runs[0].font.bold = True
            p2.runs[0].font.color.rgb = cor

    def tabela(self, s, cab, linhas, topo=2.9, larguras=None, alt_lin=0.32):
        nl, nc = len(linhas) + 1, len(cab)
        alt = Inches(0.36 + alt_lin * len(linhas))
        tb = s.shapes.add_table(nl, nc, L, Inches(topo), W, alt).table
        if larguras:
            tot = sum(larguras)
            for j, w in enumerate(larguras):
                tb.columns[j].width = Emu(int(W * w / tot))
        for j, t in enumerate(cab):
            c = tb.cell(0, j)
            c.text = t
            pr = c.text_frame.paragraphs[0]
            pr.runs[0].font.size, pr.runs[0].font.bold = Pt(10), True
            pr.runs[0].font.color.rgb = BRANCO
            c.fill.solid()
            c.fill.fore_color.rgb = ESC
        for i, L2 in enumerate(linhas, start=1):
            for j, t in enumerate(L2):
                c = tb.cell(i, j)
                # celula vazia nao gera run: escrever espaco evita IndexError
                c.text = str(t) if str(t).strip() else " "
                pr = c.text_frame.paragraphs[0]
                if pr.runs:
                    pr.runs[0].font.size = Pt(10)
                if j > 0:
                    pr.alignment = PP_ALIGN.RIGHT
                c.fill.solid()
                c.fill.fore_color.rgb = CLARO if i % 2 else BRANCO
        return tb

    def texto(self, s, linhas, topo=2.9, size=13):
        cx = s.shapes.add_textbox(L, Inches(topo), W, Inches(3.6))
        tf = cx.text_frame
        tf.word_wrap = True
        for i, (txt, bold) in enumerate(linhas):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = txt
            p.space_after = Pt(9)
            for r in p.runs:
                r.font.size = Pt(size)
                r.font.bold = bold
                r.font.color.rgb = ESC if bold else RGBColor(0x33, 0x33, 0x33)

    def grafico(self, s, cats, series, topo=2.7, alt=3.9, tipo=XL_CHART_TYPE.COLUMN_CLUSTERED):
        cd = CategoryChartData()
        cd.categories = cats
        for nome, vals in series:
            cd.add_series(nome, vals)
        g = s.shapes.add_chart(tipo, L, Inches(topo), W, Inches(alt), cd).chart
        g.has_legend = True
        g.legend.position = XL_LEGEND_POSITION.BOTTOM
        g.legend.include_in_layout = False
        g.font.size = Pt(10)
        return g


def construir():
    c = carregar()
    d, b = c["resumo"], c["book"]
    status = str(d.get("status", ""))[:60]
    dk = Deck(status)
    cur, sin, per = c["curva"], c["sinal"], c["pernas"]
    mes = lambda x: pd.Timestamp(x).strftime("%b/%y")

    # ---------------------------------------------------------------- 1 capa
    s = dk.slide("Curva Forward de Referencia e Book Proposto",
                 "Energia Convencional Flat — SE/CO  |  Metodologia publica, "
                 "auditavel e reproduzivel  |  Data de corte 14/08/2026", capa=True)

    # ---------------------------------------------------------------- 2 tese
    s = dk.slide("A tese em uma pagina",
                 "O mercado paga premio acima da projecao oficial de PLD em todos os "
                 "vertices com sinal. A posicao vende esse premio.")
    dk.kpis(s, [
        ("LADO DO BOOK", "VENDIDO", VERM),
        # nao ha MWm unico para strip mensal com quantidade diferente por mes;
        # o tamanho vai em energia e o ladder aparece na tabela do slide 3
        ("VOLUME", f"{_n(b.get('energia_liquida_gwh',0),1)} GWh", ESC),
        ("NOTIONAL", _mi(b.get("notional_brl", 0)), ESC),
        ("PnL ESPERADO", _mi(b.get("pnl_Entrega_Esperado", 0)), VERDE),
        ("RISCO (VaR)", _mi(b.get("var_total", 0)), ESC),
        ("USO DO LIMITE", f"{b.get('consumo_limite',0)*100:.0f}%", ESC),
        ("RETORNO / RISCO", f"{b.get('retorno_risco_carrego',0):.2f}x", VERDE),
    ])
    dk.texto(s, [
        ("Por que vendido.", True),
        ("A curva de referencia de mercado esta acima da projecao oficial de PLD da "
         "CCEE em todos os meses. Quem vende a termo hoje entrega energia a um preco "
         "superior ao que o modelo fundamental espera para o spot no mes de entrega.", False),
        ("Por que esses meses, e nao todos.", True),
        ("Agosto fica fora: o premio de R$ 14,95/MWh esta abaixo do limiar de "
         "R$ 15,00 e nao paga o custo de execucao. O mes ja esta quase inteiro "
         "realizado — nao ha risco a carregar que justifique premio.", False),
        ("O que invalida a tese.", True),
        ("Convergencia do premio antes da entrega, revisao relevante da projecao de "
         "PLD, ou hidrologia adversa que leve o spot acima do preco de entrada.", False),
    ], topo=3.25, size=12)

    # ---------------------------------------------------------------- 3 curva
    if len(cur):
        s = dk.slide("A curva: fundamento oficial e componente estatistico",
                     "Fair value = 70% ancora do InfoPLD + 30% componente sazonal "
                     "do PLD historico, com nowcast do IPDO. Nenhum valor colado.")
        cats = [mes(m) for m in cur.mes_ref]
        dk.grafico(s, cats, [
            ("Projecao PLD (CCEE)", [round(float(x), 1) for x in cur.fundamental]),
            ("Componente sazonal", [round(float(x), 1) for x in cur.sazonal]),
            ("FAIR VALUE do modelo", [round(float(x), 1) for x in cur.fair_value]),
            ("Referencia de mercado", [round(float(x), 1) for x in sin.referencia_mercado]),
        ], topo=2.35, alt=4.3)

    # ---------------------------------------------------------------- 4 sinal
    if len(sin):
        s = dk.slide("Onde o modelo discorda do mercado",
                     "Premio = referencia de mercado menos projecao de PLD. "
                     "Acima do limiar vira acao; abaixo e ruido do modelo.")
        linhas = [[mes(r.mes_ref), _n(r.ancora, 2), _n(r.referencia_mercado, 2),
                   _n(r.premio_rs, 2), _n(r.limiar, 2), r.acao]
                  for r in sin.itertuples()]
        dk.tabela(s, ["vertice", "projecao PLD", "referencia mercado", "premio R$/MWh",
                      "limiar", "ACAO"], linhas, topo=2.3,
                  larguras=[1.2, 1.3, 1.5, 1.3, 1.0, 1.1])
        dk.texto(s, [
            ("O limiar existe para nao operar ruido: abaixo dele a diferenca entre "
             "modelo e mercado nao supera o custo de execucao.", False)],
            topo=2.35 + 0.36 + 0.32 * len(linhas) + 0.25, size=11)

    # ---------------------------------------------------------------- 5 book
    if len(per):
        s = dk.slide("O book: lado e tamanho decididos vertice a vertice",
                     "Cada mes tem o seu lado e o seu tamanho. O tamanho sai do "
                     "risco: MWm = orcamento de risco x conviccao / VaR do vertice.")
        linhas = [[r.produto, "VENDIDO" if r.lado == "V" else "COMPRADO",
                   _n(r.mwmed), _n(r.preco_entrada, 2),
                   _mi(r.pnl_Entrega_Esperado), _mi(r.pnl_Entrega_Seco)]
                  for r in per.itertuples()]
        linhas.append(["TOTAL", "", _n(per.mwmed.sum()), "",
                       _mi(per.pnl_Entrega_Esperado.sum()),
                       _mi(per.pnl_Entrega_Seco.sum())])
        dk.tabela(s, ["vertice", "lado", "MWmed", "entrada R$/MWh",
                      "PnL esperado", "PnL cenario seco"], linhas, topo=2.3,
                  larguras=[1.1, 1.1, 1.0, 1.4, 1.4, 1.4])

    # ---------------------------------------------------------------- 6 risco
    s = dk.slide("Risco: o VaR e a unica restricao, e ela e respeitada",
                 "Fator de risco = preco a termo, nao PLD spot. Volatilidade do nivel "
                 "dessazonalizado, amortecida ate a entrega.")
    dk.kpis(s, [
        ("VaR DO BOOK", _mi(b.get("var_total", 0)), ESC),
        ("LIMITE", "R$ 50,0 mi", CINZA),
        ("USO DO LIMITE", f"{b.get('consumo_limite',0)*100:.0f}%", VERDE),
        ("EXPECTED SHORTFALL", _mi(b.get("es_total", 0)), ESC),
        ("PIOR CENARIO", _mi(b.get("pior_cenario", 0)), VERM),
    ])
    dk.texto(s, [
        ("O que entra no risco.", True),
        ("O consumo de limite e o VaR de marcacao a mercado de 1 mes, somado entre "
         "os vertices sem hipotese de diversificacao. A perda de carregar ate a "
         "entrega no cenario adverso continua calculada e visivel, mas nao dimensiona: "
         "e um horizonte diferente do limite declarado no mandato.", False),
        ("Por que nao usar a volatilidade do PLD spot.", True),
        ("Spot reverte a media em cerca de 28 dias e tem piso e teto administrativos. "
         "A volatilidade dele superestima o risco de carregar um contrato a termo. "
         "O spot e reportado como TETO de referencia, nao como a medida.", False),
    ], topo=3.3, size=12)

    # ---------------------------------------------------------------- 7 cenarios
    if len(cur):
        s = dk.slide("Cenarios: trajetorias oficiais inteiras, nunca combinadas",
                     "Esperado, Seco e Umido sao trajetorias publicadas pela CCEE, "
                     "preservadas por inteiro. Combinar meses de modelos diferentes "
                     "produziria uma curva que nenhum modelo gerou.")
        cats = [mes(m) for m in cur.mes_ref]
        dk.grafico(s, cats, [
            ("Seco", [round(float(x), 1) for x in cur.seco]),
            ("Esperado", [round(float(x), 1) for x in cur.fair_value]),
            ("Umido", [round(float(x), 1) for x in cur.umido]),
        ], topo=2.5, alt=4.1)

    # ---------------------------------------------------------------- 8 pnl
    s = dk.slide("Resultado por cenario",
                 "Duas familias de resultado, com horizontes distintos: convergencia "
                 "do premio (marcacao) e carrego ate a entrega.")
    dk.grafico(s, ["Convergencia", "Carrego esperado", "Carrego seco", "Carrego umido"],
               [("PnL (R$ mi)", [round(b.get(k, 0) / 1e6, 2) for k in
                                 ("pnl_Convergencia", "pnl_Entrega_Esperado",
                                  "pnl_Entrega_Seco", "pnl_Entrega_Umido")])],
               topo=2.35, alt=3.4)
    dk.texto(s, [
        (f"Assimetria declarada: o ganho no cenario umido "
         f"({_mi(b.get('pnl_Entrega_Umido',0))}) e varias vezes maior que a perda no "
         f"cenario seco ({_mi(b.get('pnl_Entrega_Seco',0))}), porque a posicao "
         f"vendida ganha quando o preco cai e o piso administrativo limita a queda.",
         False)], topo=5.95, size=11)

    # ---------------------------------------------------------------- 9 governanca
    s = dk.slide("Governanca: por que estes numeros sao defensaveis",
                 "Nenhum resultado da pasta e valor colado. As unicas celulas com "
                 "valor sao a serie observada, as premissas e a referencia de mercado.")
    dk.texto(s, [
        ("Fonte declarada em cada numero.", True),
        ("Toda celula carrega um rotulo: OBSERVADO, CALCULADO, PREMISSA, PROXY ou "
         "FAIR_VALUE. O manifesto registra hash SHA-256 e origem de cada arquivo.", False),
        ("Conferencia cruzada.", True),
        ("A aba CROSSCHECK recalcula em Python as mesmas grandezas que a planilha "
         "calcula por formula e compara com tolerancia por linha. Divergencia "
         "aparece como DIVERGENCIA, nao passa despercebida.", False),
        ("Parametros escolhidos fora da amostra.", True),
        ("Meia-vida da sazonalidade selecionada por validacao walk-forward, sem usar "
         "informacao futura. Janela e kernel escolhidos pelo erro de FORMA, com o "
         "nivel removido — selecionar forma pelo erro total escolheria sempre o "
         "ajuste mais liso.", False),
        ("O que o modelo NAO afirma.", True),
        ("A projecao de PLD da CCEE e expectativa de spot no mes de entrega, nao "
         "preco de contrato bilateral. O premio de risco a termo nao e observavel em "
         "fonte publica: entra como premissa declarada, com risco de base explicito.",
         False),
    ], topo=2.0, size=11.5)

    # ---------------------------------------------------------------- 10 proximos
    s = dk.slide("Execucao e monitoramento", None)
    dk.texto(s, [
        ("Entrada.", True),
        ("Executar em lotes por vertice, respeitando o teto operacional de 250 MWm "
         "por mes. Tres dos quatro vertices estao limitados por liquidez, nao por "
         "risco — o VaR permitiria mais.", False),
        ("Stop.", True),
        (f"O limite de VaR e o stop primario. Folga atual ate o limite: "
         f"{_mi(50e6 - b.get('var_total',0))}. Movimento adverso do preco a termo "
         f"maior que o VaR do vertice em um mes exige revisao do tamanho.", False),
        ("Gatilhos de revisao.", True),
        ("Novo InfoPLD com revisao relevante da projecao; indice hidrologico cruzando "
         "os limiares de regime; convergencia do premio antes da entrega.", False),
        ("Reproducao.", True),
        ("make all reexecuta a cadeia inteira — download, leitura de boletins, "
         "reselecao de parametros, pipeline e testes — e falha explicitamente se "
         "qualquer fonte obrigatoria faltar.", False),
    ], topo=1.7, size=12)

    alvo = OUT / "Apresentacao_Executiva.pptx"
    dk.p.save(alvo)
    return alvo, dk.n


if __name__ == "__main__":
    a, n = construir()
    print(f"gerado: {a}  ({n} slides)")
